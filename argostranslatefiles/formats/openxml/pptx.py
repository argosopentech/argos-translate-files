from argostranslate.translate import ITranslation
from pptx import Presentation
from pptx.chart.data import CategoryChartData,BubbleChartData,XyChartData, XySeriesData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

import sys
from argostranslatefiles.formats.abstract_xml import AbstractXml


class Pptx(AbstractXml):
    supported_file_extensions = ['.pptx']
    bubble_chart_types = [XL_CHART_TYPE.BUBBLE, XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT]
    xy_chart_types = [XL_CHART_TYPE.XY_SCATTER,XL_CHART_TYPE.XY_SCATTER_LINES,XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,XL_CHART_TYPE.XY_SCATTER_SMOOTH,XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS]

    def translate_paragraphs(self, paras, underlying_translation):
        for paragraph in paras:
            #get the average font size for the para
            if paragraph.runs:
                font_sizes = [(x.font.size,len(x.text)) for x in paragraph.runs]
                #sys.stderr.write(f"paragraph font size: {paragraph.font.size}. font sizes: {str(font_sizes)}\n")
                if font_sizes:
                    #avg_font_size = round(sum([x[0] for x in font_sizes])/len([x[0] for x in font_sizes]))
                    most_common_font_size = sorted(font_sizes,key=lambda x: x[1])[0]
                    if most_common_font_size:
                        paragraph.font.size = most_common_font_size[0]
                paragraph.text = underlying_translation.translate(paragraph.text)

    def process_shape(self, shape, underlying_translation):
        if shape.has_text_frame:
            self.translate_paragraphs(shape.text_frame.paragraphs, underlying_translation)
    
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    self.translate_paragraphs(cell.text_frame.paragraphs, underlying_translation)
        if shape.has_chart:
            chart = shape.chart
        
            if chart.has_title: 
                if chart.chart_title.has_text_frame:
                    self.translate_paragraphs(chart.chart_title.text_frame.paragraphs, underlying_translation)
            # TODO: Bubble and Xy charts are parsed differently, not clear how to retrieve values
            if chart.chart_type in self.bubble_chart_types:
                new_chart_data = BubbleChartData()
                return
            elif chart.chart_type in self.xy_chart_types:
                new_chart_data = XyChartData()
                return
            else:
                new_chart_data = CategoryChartData()
            new_categories = [underlying_translation.translate(c) for c in chart.plots[0].categories]
            new_chart_data.categories = new_categories
            for series in chart.series:
                new_chart_data.add_series(underlying_translation.translate(series.name),series.values)
            chart.replace_data(new_chart_data)

            # TODO: Add SmartArt handling here. Pypptx does not support SmartArt, but it might be
            # possible to first convert SmartArt to shapes by some means. Another alternative is just
            # to parse the xml to find things that look like text

    def translate(self, underlying_translation: ITranslation, file_path: str):
        outzip_path = self.get_output_path(underlying_translation, file_path)
        presentation = Presentation(file_path)
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for subshape in shape.shapes:
                        self.process_shape(subshape, underlying_translation)
                else:
                    self.process_shape(shape, underlying_translation)

        presentation.save(outzip_path)  


        return outzip_path
